from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from io import BytesIO
import openpyxl
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
import os
import locale
import traceback
import base64
import tempfile
import subprocess
import pythoncom
import win32com.client as win32
from pdf2image import convert_from_path
from PIL import Image
import json

# ===================================================================
# CONFIGURAÇÕES
# ===================================================================

app = Flask(__name__)
CORS(app)

try:
    locale.setlocale(locale.LC_ALL, 'pt_BR.UTF-8')
except locale.Error:
    print("Aviso: localidade pt_BR.UTF-8 não encontrada. Usando fallback.")

# Caminhos externos - AJUSTE ESTES CAMINHOS PARA O SEU AMBIENTE
POPPLER_PATH = r"C:\Users\fixxi\Desktop\STYLUX_PPTGEN\Back\POPPLER\Library\bin"
LIBREOFFICE_PATH = r"C:\Program Files\LibreOffice\program\soffice.exe"

# ===================================================================
# FUNÇÕES DE APOIO
# ===================================================================

def format_cell(cell):
    if cell.value is None:
        return ""
    if '%' in str(cell.number_format):
        try:
            return f"{float(cell.value):.2%}".replace('.', ',')
        except (ValueError, TypeError):
            return str(cell.value)
    if 'R$' in str(cell.number_format) or 'BRL' in str(cell.number_format):
        try:
            return locale.currency(float(cell.value), grouping=True)
        except (ValueError, TypeError):
            return str(cell.value)
    if isinstance(cell.value, (float, int)):
        return f"{cell.value:,.2f}".replace(',', 'X').replace('.', ',').replace('X', '.')
    return str(cell.value)

def build_table_data(ws_formats, ws_values, range_string):
    table_data = []
    for row_format, row_value in zip(ws_formats[range_string], ws_values[range_string]):
        new_row = []
        for cell_format, cell_value in zip(row_format, row_value):
            cell_format.value = cell_value.value
            new_row.append(format_cell(cell_format))
        table_data.append(new_row)
    return table_data


def substituir_textos(prs, substituicoes, campos_ativos):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame or shape.has_table:
                continue
            for paragraph in shape.text_frame.paragraphs:
                runs_data = [(run.text, run.font) for run in paragraph.runs]
                full_text = ''.join(run.text for run in paragraph.runs)
                original_text = full_text

                for chave, valor in substituicoes.items():
                    nome_campo = chave.strip('{} ')
                    if nome_campo in campos_ativos and chave in full_text:
                        full_text = full_text.replace(chave, str(valor) if valor is not None else "")

                placeholders_inativos = {chave for chave in substituicoes if chave.strip('{} ') not in campos_ativos}
                for chave_inativa in placeholders_inativos:
                    if chave_inativa in full_text:
                        full_text = full_text.replace(chave_inativa, "")

                if original_text != full_text:
                    for _ in range(len(paragraph.runs)):
                        p = paragraph.runs[0]._r
                        p.getparent().remove(p)
                    if runs_data:
                        new_run = paragraph.add_run()
                        new_run.text = full_text
                        original_font = runs_data[0][1]
                        new_run.font.name = original_font.name
                        new_run.font.size = original_font.size
                        new_run.font.bold = original_font.bold
                        new_run.font.italic = original_font.italic
                        new_run.font.underline = original_font.underline
                        original_color = original_font.color
                        if original_color.type == MSO_COLOR_TYPE.RGB:
                            new_run.font.color.rgb = original_color.rgb
                        elif original_color.type == MSO_COLOR_TYPE.SCHEME:
                            new_run.font.color.theme_color = original_color.theme_color
                            new_run.font.color.brightness = original_color.brightness
                    else:
                        paragraph.add_run().text = full_text

def substituir_logo(prs, logo_stream, placeholder):
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if not shape.has_text_frame:
                continue
            if placeholder in shape.text_frame.text:
                logo_stream.seek(0)
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                sp = shape._sp
                sp.getparent().remove(sp)
                slide.shapes.add_picture(logo_stream, left, top, width, height)

def substituir_tabela(slide, placeholder, table_data, campos_ativos):
    if placeholder.strip('{} ') not in campos_ativos:
        for shape in slide.shapes:
            if shape.has_text_frame and placeholder.lower() in shape.text_frame.text.lower():
                sp = shape._sp
                sp.getparent().remove(sp)
        return

    shape_to_replace = None
    for shape in slide.shapes:
        if shape.has_text_frame and placeholder.lower() in shape.text_frame.text.lower():
            shape_to_replace = shape
            break
    if shape_to_replace:
        left, top, width, height = shape_to_replace.left, shape_to_replace.top, shape_to_replace.width, shape_to_replace.height

        sp = shape_to_replace._sp
        sp.getparent().remove(sp)

        num_rows = len(table_data)
        num_cols = len(table_data[0]) if table_data else 0
        if num_rows == 0 or num_cols == 0:
            return

        table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        tbl = table_shape.table
        for r_idx, row_data in enumerate(table_data):
            for c_idx, cell_data in enumerate(row_data):
                cell = tbl.cell(r_idx, c_idx)
                cell.text = str(cell_data)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(15)
                para.font.name = 'Calibri'
                para.font.color.rgb = RGBColor(0, 0, 0)
                if r_idx == 0:
                    para.font.bold = True
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(220, 220, 220)


def substituir_graficos(prs, excel_path, graficos_info):
    """
    Substitui placeholders de gráficos no PPT com imagens do Excel de forma robusta.
    """
    print("--- Iniciando substituição de gráficos ---")
    excel = None
    wb = None
    
    try:
        pythoncom.CoInitialize()
        excel = win32.gencache.EnsureDispatch("Excel.Application")
        excel.Visible = False
        excel.DisplayAlerts = False

        print(f"Abrindo workbook: {excel_path}")
        wb = excel.Workbooks.Open(excel_path)
        
        print("Atualizando dados e gráficos no Excel...")
        wb.RefreshAll()
        excel.CalculateUntilAsyncQueriesDone()
        
        # Otimização: Mapeia todos os gráficos do Excel primeiro
        charts_map = {}
        for ws in wb.Worksheets:
            try:
                for chart_object in ws.ChartObjects():
                    charts_map[chart_object.Name] = chart_object
            except Exception:
                # Planilha pode não ter ChartObjects, ignora o erro
                pass
        print(f"Gráficos encontrados no Excel: {list(charts_map.keys())}")

        # Itera pelos slides para encontrar e substituir placeholders
        for slide_idx, slide in enumerate(prs.slides):
            # Usar list() para criar uma cópia, permitindo remover shapes durante a iteração
            for shape in list(slide.shapes):
                if not shape.has_text_frame:
                    continue
                
                # Verifica se o texto do shape corresponde a algum placeholder de gráfico
                for placeholder, chart_name in graficos_info.items():
                    if placeholder in shape.text_frame.text:
                        print(f"-> Placeholder '{placeholder}' encontrado no Slide {slide_idx + 1}.")
                        
                        if chart_name in charts_map:
                            print(f"--> Gráfico '{chart_name}' correspondente encontrado no Excel.")
                            
                            # Exporta o gráfico para um arquivo temporário
                            chart_object = charts_map[chart_name]
                            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_file:
                                image_path = tmp_file.name
                            
                            chart_object.Chart.Export(image_path)
                            print(f"--> Gráfico exportado para: {image_path}")

                            # Adiciona a imagem no lugar do placeholder
                            left, top, width, height = shape.left, shape.top, shape.width, shape.height
                            slide.shapes.add_picture(image_path, left, top, width=width, height=height)
                            
                            # Remove o shape do placeholder
                            sp = shape._sp
                            sp.getparent().remove(sp)
                            print("--> Imagem inserida e placeholder removido.")
                            
                            os.remove(image_path)
                        else:
                            print(f"AVISO: O gráfico '{chart_name}' não foi encontrado no arquivo Excel. Verifique o nome.")
                        # Sai do loop de placeholders, pois este shape já foi processado
                        break
                        
    except Exception as e:
        print(f"ERRO CRÍTICO na função substituir_graficos: {e}")
        traceback.print_exc()
    finally:
        # Garante que o Excel seja fechado corretamente
        if wb:
            wb.Close(SaveChanges=False)
        if excel:
            excel.Quit()
        pythoncom.CoUninitialize()
        print("--- Finalizada a substituição de gráficos ---")


def adicionar_slides_customizados(prs, custom_slides_json):
    try:
        custom_slides_data = json.loads(custom_slides_json)
        if not custom_slides_data:
            return

        template_slide_layout = prs.slides[-1].slide_layout

        for slide_data in custom_slides_data:
            slide_list = list(prs.slides)
            target_idx = len(slide_list) - 1
            
            new_slide = prs.slides.add_slide(template_slide_layout)

            title_text = slide_data.get('title', '')
            content_text = slide_data.get('content', '')

            if new_slide.shapes.title:
                new_slide.shapes.title.text = title_text
            
            for shape in new_slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.type == 'BODY':
                    shape.text_frame.text = content_text
                    break
            
            sldIdLst = prs.part.package.part_related_by(prs.part.reltype).target_part.element.sldIdLst
            slides = list(sldIdLst)
            last_slide_id = slides[-1]
            sldIdLst.remove(last_slide_id)
            sldIdLst.insert(target_idx, last_slide_id)

    except Exception as e:
        print(f"Erro ao adicionar slides customizados: {e}")


def create_ppt(text_substitutions, table_data1, table_data2, campos_ativos, excel_path, slides_a_manter=None, logo_stream=None, custom_slides_json=None):
    modelo_path = os.path.join(os.path.dirname(__file__), "templates_ppt", "modeloprincipal3.pptx")
    prs = Presentation(modelo_path)

    if custom_slides_json:
        adicionar_slides_customizados(prs, custom_slides_json)

    if slides_a_manter:
        indices_a_manter = {s - 1 for s in slides_a_manter}
        # Itera de trás para frente para remover slides sem afetar os índices dos slides restantes
        for i in range(len(prs.slides) - 1, -1, -1):
            if i not in indices_a_manter:
                rId = prs.slides._sldIdLst[i].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[i]
            
    campos_ativos.extend([
        "FABRICANTEMODULO", "MODELOMODULO", "POTENCIAMODULO",
        "FABRICANTEINVERSORES", "MODELOINVERSORES", "POTENCIAINVERSORES",
        "ESTRUTURA", "SISTEMAMONITORAMENTO",
        "EQUIPAMENTOSESTRUTURA", "VIDAUTIL", "CERTIFICACAO", "ISOLANTE",
        "CONTATO", "TENSAO", "PROTECAO", "TEMPOPERACAO",
        "Manual1", "Manual2", "Manual3", "Manual4", "SOMATOTALFINAL",
        "FLUXO1", "FLUXO2"
    ])

    substituir_textos(prs, text_substitutions, campos_ativos)
    
    if logo_stream:
        substituir_logo(prs, logo_stream, '{{LOGOCLIENTE}}')

    for slide in prs.slides:
        substituir_tabela(slide, '{{FLUXO1}}', table_data1, campos_ativos)
        substituir_tabela(slide, '{{FLUXO2}}', table_data2, campos_ativos)

    graficos_info = {
        "{{grafico_receita}}": "ReceitaAnual",
        "{{grafico_custos}}": "Custo"
    }
    substituir_graficos(prs, excel_path, graficos_info)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


def pptx_to_pdf(pptx_path, output_dir):
    try:
        subprocess.run(
            [
                LIBREOFFICE_PATH, "--headless", "--convert-to", "pdf",
                "--outdir", output_dir, pptx_path
            ],
            check=True,
            timeout=60
        )
        pdf_filename = os.path.basename(pptx_path).replace(".pptx", ".pdf")
        return os.path.join(output_dir, pdf_filename)
    except subprocess.TimeoutExpired:
        raise RuntimeError("A conversão para PDF demorou demais (timeout).")
    except Exception as e:
        raise RuntimeError(f"Erro ao converter PPTX para PDF: {e}")

# ===================================================================
# ENDPOINTS
# ===================================================================

@app.route("/extract", methods=["POST"])
def extract_data():
    try:
        if "file" not in request.files:
            return jsonify({"error": "Nenhum arquivo enviado."}), 400
        
        file_bytes = request.files['file'].read()

        wb_values = openpyxl.load_workbook(filename=BytesIO(file_bytes), data_only=True)
        wb_formats = openpyxl.load_workbook(filename=BytesIO(file_bytes), data_only=False)
        
        ws_values = wb_values["Extract"]
        ws_formats = wb_formats["Extract"]

        row_cells_formats = ws_formats[2]
        row_cells_values = ws_values[2]

        def get_formatted_value(index):
            if len(row_cells_formats) > index and len(row_cells_values) > index:
                cell_with_format = row_cells_formats[index]
                cell_with_value = row_cells_values[index]
                
                cell_with_format.value = cell_with_value.value
                return format_cell(cell_with_format)
            return ""

        extracted_data = {
            "NOME_CLIENTE": get_formatted_value(0),
            "CONS_ENERGIA_MEDIO": get_formatted_value(1),
            "VOL_PROJ": get_formatted_value(2),
            "OBJ": get_formatted_value(3),
            "PRAZO_CONT": get_formatted_value(4),
            "DESC_1ANO": get_formatted_value(5),
            "MODELO_NEGOCIO": get_formatted_value(6),
            "TAXA_MEDIA": get_formatted_value(12),
            "PIS": get_formatted_value(13),
            "ICMS": get_formatted_value(14),
            "PONTA": get_formatted_value(15),
            "FORA_PONTA": get_formatted_value(16),
        }
        return jsonify(extracted_data)

    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": f"Erro na extração: {e}"}), 500


@app.route("/generate", methods=["POST"])
def generate_ppt():
    excel_path = None
    try:
        if "file" not in request.files:
            return jsonify({"error": "Nenhum arquivo enviado."}), 400
            
        file_bytes = request.files["file"].read()
        logo_file = request.files.get("logo")
        logo_stream = BytesIO(logo_file.read()) if logo_file else None
        custom_slides_json = request.form.get("custom_slides")

        with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp_excel:
            tmp_excel.write(file_bytes)
            excel_path = tmp_excel.name

        wb_values = openpyxl.load_workbook(filename=BytesIO(file_bytes), data_only=True)
        wb_formats = openpyxl.load_workbook(filename=BytesIO(file_bytes), data_only=False)

        ws_values_for_text = wb_values["Extract"]
        row_values = [cell.value for cell in ws_values_for_text[2]]
        
        text_keys = [
            "NOME_CLIENTE", "CONS_ENERGIA_MEDIO", "VOL_PROJ", "OBJ", "PRAZO_CONT",
            "DESC_1ANO", "MODELO_NEGOCIO", "", "", "", "", "", "TAXA_MEDIA", "PIS",
            "ICMS", "PONTA", "FORA_PONTA"
        ]
        text_subs = {
            f"{{{{{key}}}}}": (str(row_values[idx]) if idx < len(row_values) and row_values[idx] is not None else "")
            for idx, key in enumerate(text_keys) if key
        }
        
        form_data = request.form.to_dict()
        for key, value in form_data.items():
            if key not in ["campos", "slides_a_manter", "custom_slides"]:
                text_subs[f"{{{{{key}}}}}"] = value

        ws_formats_for_table = wb_formats["Extract"]
        ws_values_for_table = wb_values["Extract"]
        table_data1 = build_table_data(ws_formats_for_table, ws_values_for_table, 'H2:L17')
        table_data2 = build_table_data(ws_formats_for_table, ws_values_for_table, 'R2:V17')

        campos_ativos = request.form.getlist("campos")
        slides_a_manter = request.form.getlist("slides_a_manter", type=int)
        
        ppt_buffer = create_ppt(text_subs, table_data1, table_data2, campos_ativos, excel_path, slides_a_manter, logo_stream=logo_stream, custom_slides_json=custom_slides_json)

        return send_file(
            ppt_buffer,
            as_attachment=True,
            download_name="proposta_customizada.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": f"Erro na geração: {e}"}), 500
    finally:
        if excel_path and os.path.exists(excel_path):
            os.remove(excel_path)


@app.route("/preview", methods=["POST"])
def preview_ppt():
    excel_path = None
    pptx_temp_file = None
    pdf_temp_file = None
    
    try:
        if "file" not in request.files:
            return jsonify({"error": "Nenhum arquivo enviado."}), 400
            
        file_bytes = request.files["file"].read()
        logo_file = request.files.get("logo")
        logo_stream = BytesIO(logo_file.read()) if logo_file else None
        custom_slides_json = request.form.get("custom_slides")
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp_excel:
            tmp_excel.write(file_bytes)
            excel_path = tmp_excel.name

        wb_values = openpyxl.load_workbook(filename=BytesIO(file_bytes), data_only=True)
        wb_formats = openpyxl.load_workbook(filename=BytesIO(file_bytes), data_only=False)

        ws_values_for_text = wb_values["Extract"]
        row_values = [cell.value for cell in ws_values_for_text[2]]

        text_keys = [
            "NOME_CLIENTE", "CONS_ENERGIA_MEDIO", "VOL_PROJ", "OBJ", "PRAZO_CONT",
            "DESC_1ANO", "MODELO_NEGOCIO", "", "", "", "", "", "TAXA_MEDIA", "PIS",
            "ICMS", "PONTA", "FORA_PONTA"
        ]
        text_subs = {
            f"{{{{{key}}}}}": (str(row_values[idx]) if idx < len(row_values) and row_values[idx] is not None else "")
            for idx, key in enumerate(text_keys) if key
        }
        
        form_data = request.form.to_dict()
        for key, value in form_data.items():
            if key not in ["campos", "slides_a_manter", "custom_slides"]:
                text_subs[f"{{{{{key}}}}}"] = value
        
        ws_formats_for_table = wb_formats["Extract"]
        ws_values_for_table = wb_values["Extract"]
        table_data1 = build_table_data(ws_formats_for_table, ws_values_for_table, 'H2:L17')
        table_data2 = build_table_data(ws_formats_for_table, ws_values_for_table, 'R2:V17')
        
        campos_ativos = request.form.getlist("campos")
        slides_a_manter = request.form.getlist("slides_a_manter", type=int)
        
        ppt_buffer = create_ppt(text_subs, table_data1, table_data2, campos_ativos, excel_path, slides_a_manter, logo_stream=logo_stream, custom_slides_json=custom_slides_json)
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_pptx:
            tmp_pptx.write(ppt_buffer.getvalue())
            pptx_temp_file = tmp_pptx.name
            
        output_dir = os.path.dirname(pptx_temp_file)
        pdf_temp_file = pptx_to_pdf(pptx_temp_file, output_dir)
        
        images = convert_from_path(pdf_temp_file, poppler_path=POPPLER_PATH)
        base64_images = []
        
        for img in images:
            buffered = BytesIO()
            img.save(buffered, format="PNG")
            base64_images.append(base64.b64encode(buffered.getvalue()).decode("utf-8"))

        return jsonify({"slides": base64_images})
    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": f"Erro no preview: {e}"}), 500
    finally:
        # Limpeza robusta de arquivos temporários
        for temp_file in [pptx_temp_file, pdf_temp_file, excel_path]:
            if temp_file and os.path.exists(temp_file):
                try:
                    os.remove(temp_file)
                except Exception as e:
                    print(f"Erro ao remover arquivo temporário {temp_file}: {e}")

# ===================================================================
# MAIN
# ===================================================================

if __name__ == "__main__":
    app.run(debug=True, port=5000)